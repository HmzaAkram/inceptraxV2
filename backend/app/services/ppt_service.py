"""
PPT Service — Multi-theme 10-slide investor deck generator.

Shows only real user data. Missing stages are silently skipped/minimised.
"""
import os
import json
import re
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── Theme Configurations ──────────────────────────────────────────────────────

THEMES = {
    "dark_executive": {
        "bg":        "0A0A14",
        "card":      "111827",
        "accent":    "6366F1",
        "accent2":   "8B5CF6",
        "text":      "F1F5F9",
        "muted":     "94A3B8",
        "border":    "1E2A3A",
        "success":   "10B981",
        "warning":   "F59E0B",
        "danger":    "EF4444",
        "cover_bg":  "0A0A14",
        "name":      "Dark Executive",
    },
    "clean_light": {
        "bg":        "FFFFFF",
        "card":      "F8FAFC",
        "accent":    "6366F1",
        "accent2":   "4F46E5",
        "text":      "0F172A",
        "muted":     "64748B",
        "border":    "E2E8F0",
        "success":   "059669",
        "warning":   "D97706",
        "danger":    "DC2626",
        "cover_bg":  "0F172A",
        "name":      "Clean Light",
    },
    "gradient_bold": {
        "bg":        "0D1117",
        "card":      "161B22",
        "accent":    "00D4FF",
        "accent2":   "A855F7",
        "text":      "FFFFFF",
        "muted":     "8B949E",
        "border":    "21262D",
        "success":   "3FB950",
        "warning":   "D29922",
        "danger":    "F85149",
        "cover_bg":  "0D1117",
        "name":      "Gradient Bold",
    },
}

# ─── Helper Utilities ──────────────────────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_slide_bg(slide, hex_color: str):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex_color)


def _add_rect(slide, x, y, w, h, fill_hex=None, border_hex=None, border_pt=0.5):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = Pt(border_pt) if border_hex else 0
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
    else:
        shape.fill.background()
    if border_hex:
        shape.line.color.rgb = _rgb(border_hex)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text: str, x, y, w, h,
              font_size=11, bold=False, color_hex="F1F5F9",
              align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color_hex)
    return txBox


def _add_top_bar(slide, accent_hex: str, slide_w_in=10.0):
    bar = slide.shapes.add_shape(1, 0, 0, Inches(slide_w_in), Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent_hex)
    bar.line.fill.background()


def _add_slide_number(slide, num: int, theme: dict):
    _add_text(slide, str(num), 9.5, 5.3, 0.4, 0.2,
              font_size=8, color_hex=theme["muted"], align=PP_ALIGN.RIGHT)


def _slide_title(slide, text: str, theme: dict, y=0.18):
    _add_text(slide, text, 0.5, y, 9.0, 0.45,
              font_size=22, bold=True, color_hex=theme["text"])


def _score_color(score, theme):
    if score >= 75:
        return theme["success"]
    if score >= 50:
        return theme["warning"]
    return theme["danger"]


def _parse_list(raw, max_items=6):
    """Safely extract a list from a raw value (list, str, or None). Returns list of strings."""
    if not raw:
        return []
    if isinstance(raw, list):
        return [str(i) for i in raw[:max_items] if i]
    if isinstance(raw, str):
        lines = [l.strip().lstrip("•-*") for l in raw.split("\n") if l.strip()]
        return lines[:max_items]
    return []


def _parse_str(raw, fallback="") -> str:
    if not raw:
        return fallback
    if isinstance(raw, (dict, list)):
        return ""
    return str(raw).strip()


def _get_stage(analysis_data: dict, stage_name: str) -> dict:
    """Return the parsed JSON of a completed stage, or empty dict."""
    stages = analysis_data.get("stages", {})
    raw = stages.get(stage_name)
    if not raw:
        return {}
    if isinstance(raw, dict):
        return raw
    if isinstance(raw, str):
        try:
            return json.loads(raw)
        except Exception:
            return {}
    return {}


# ─── Slide Builder Functions ───────────────────────────────────────────────────

def build_cover_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["cover_bg"])
    _add_top_bar(slide, theme["accent"])

    title = _parse_str(data.get("title"), "Untitled Idea")
    one_liner = _parse_str(data.get("one_liner") or data.get("description", "")[:120])
    industry = _parse_str(data.get("industry") or data.get("market", ""))
    score = data.get("overall_score") or 0

    # Decorative block top-right
    _add_rect(slide, 8.2, 0.1, 1.5, 1.5, fill_hex=theme["accent"], border_hex=None)
    _add_rect(slide, 8.7, 0.5, 1.0, 1.0, fill_hex=theme["accent2"], border_hex=None)

    # Title
    _add_text(slide, title, 0.5, 0.9, 8.0, 1.5,
              font_size=44, bold=True, color_hex=theme["text"], wrap=True)

    # One-liner
    if one_liner:
        _add_text(slide, one_liner, 0.5, 2.5, 7.5, 0.7,
                  font_size=14, color_hex=theme["muted"], wrap=True)

    # Score badge + Industry
    badge_x = 0.5
    if score and score > 0:
        score_col = _score_color(score, theme)
        badge = _add_rect(slide, badge_x, 3.4, 1.1, 0.55, fill_hex=score_col)
        _add_text(slide, f"{score}", badge_x + 0.05, 3.43, 1.0, 0.5,
                  font_size=18, bold=True, color_hex="FFFFFF", align=PP_ALIGN.CENTER)
        badge_x += 1.3

    if industry:
        _add_rect(slide, badge_x, 3.4, 1.6, 0.55, fill_hex=theme["card"], border_hex=theme["accent"])
        _add_text(slide, industry, badge_x + 0.08, 3.43, 1.45, 0.5,
                  font_size=11, color_hex=theme["accent"], bold=True)

    # Branding footer
    _add_text(slide, "Powered by Inceptrax", 7.0, 5.1, 2.8, 0.3,
              font_size=8, color_hex=theme["muted"], align=PP_ALIGN.RIGHT)


def build_problem_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 2, theme)
    _slide_title(slide, "The Problem", theme)

    v = _get_stage(data, "validation")
    if not v:
        return

    problem = _parse_str(v.get("problem") or v.get("problem_statement"))
    risks = _parse_list(v.get("risks") or v.get("challenges"), 3)

    if problem:
        # Highlighted problem card
        _add_rect(slide, 0.5, 0.75, 9.0, 1.1, fill_hex=theme["card"], border_hex=theme["accent"])
        # Accent left border
        _add_rect(slide, 0.5, 0.75, 0.06, 1.1, fill_hex=theme["accent"])
        _add_text(slide, problem, 0.7, 0.82, 8.6, 1.0,
                  font_size=12, color_hex=theme["text"], wrap=True)

    # Pain point cards
    pain_points = _parse_list(v.get("pain_points") or risks, 3)
    for i, pt in enumerate(pain_points):
        cx = 0.5 + i * 3.1
        _add_rect(slide, cx, 2.1, 2.9, 2.5, fill_hex=theme["card"], border_hex=theme["border"])
        # Number circle
        _add_rect(slide, cx + 0.1, 2.2, 0.5, 0.45, fill_hex=theme["accent"])
        _add_text(slide, str(i + 1), cx + 0.15, 2.22, 0.4, 0.4,
                  font_size=12, bold=True, color_hex="FFFFFF", align=PP_ALIGN.CENTER)
        _add_text(slide, pt[:180], cx + 0.15, 2.75, 2.65, 1.7,
                  font_size=10, color_hex=theme["text"], wrap=True)


def build_solution_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 3, theme)
    _slide_title(slide, "Our Solution", theme)

    v = _get_stage(data, "validation")
    if not v:
        return

    solution = _parse_str(v.get("solution") or v.get("solution_description"))
    strengths = _parse_list(v.get("strengths") or v.get("key_strengths"), 4)

    if solution:
        _add_rect(slide, 0.5, 0.75, 5.5, 1.3, fill_hex=theme["card"], border_hex=theme["border"])
        _add_text(slide, solution, 0.65, 0.82, 5.2, 1.2,
                  font_size=10, color_hex=theme["text"], wrap=True)

    # Steps / how it works
    how_it_works = _parse_list(v.get("how_it_works") or v.get("key_features"), 4)
    steps = how_it_works or strengths
    for i, step in enumerate(steps[:4]):
        sy = 2.2 + i * 0.65
        _add_rect(slide, 0.5, sy, 0.38, 0.38, fill_hex=theme["accent"])
        _add_text(slide, str(i + 1), 0.52, sy + 0.01, 0.36, 0.36,
                  font_size=10, bold=True, color_hex="FFFFFF", align=PP_ALIGN.CENTER)
        _add_text(slide, step[:120], 1.05, sy, 4.8, 0.5,
                  font_size=10, color_hex=theme["text"])

    # Strength cards right side
    strength_items = _parse_list(v.get("strengths"), 3)
    for i, s in enumerate(strength_items[:3]):
        sy = 0.75 + i * 1.55
        _add_rect(slide, 6.2, sy, 3.3, 1.35, fill_hex=theme["card"], border_hex=theme["success"])
        _add_text(slide, "✓", 6.35, sy + 0.05, 0.4, 0.4,
                  font_size=13, bold=True, color_hex=theme["success"])
        _add_text(slide, s[:130], 6.75, sy + 0.08, 2.7, 1.1,
                  font_size=10, color_hex=theme["text"], wrap=True)


def build_market_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 4, theme)
    _slide_title(slide, "Market Opportunity", theme)

    mr = _get_stage(data, "market_research")
    if not mr:
        return

    tam = _parse_str(mr.get("tam") or mr.get("total_addressable_market"))
    sam = _parse_str(mr.get("sam") or mr.get("serviceable_addressable_market"))
    som = _parse_str(mr.get("som") or mr.get("serviceable_obtainable_market"))
    cagr = _parse_str(mr.get("cagr") or mr.get("growth_rate"))
    trends = _parse_list(mr.get("trends") or mr.get("market_trends"), 4)

    # TAM / SAM / SOM boxes
    metrics = [("TAM", tam), ("SAM", sam), ("SOM", som)]
    for i, (label, val) in enumerate(metrics):
        if not val:
            continue
        mx = 0.5 + i * 3.1
        _add_rect(slide, mx, 0.75, 2.9, 1.3, fill_hex=theme["card"], border_hex=theme["accent"])
        _add_text(slide, label, mx + 0.15, 0.82, 2.6, 0.4,
                  font_size=11, bold=True, color_hex=theme["muted"])
        _add_text(slide, val[:30], mx + 0.15, 1.12, 2.6, 0.6,
                  font_size=16, bold=True, color_hex=theme["accent"])

    # CAGR highlight
    if cagr:
        _add_rect(slide, 0.5, 2.25, 3.0, 0.65, fill_hex=theme["accent"], border_hex=None)
        _add_text(slide, f"CAGR: {cagr}", 0.65, 2.32, 2.7, 0.5,
                  font_size=14, bold=True, color_hex="FFFFFF")

    # Simple visual bar for market trends (fake proportional bars)
    if trends:
        _add_text(slide, "Key Market Trends", 0.5, 3.1, 5.0, 0.35,
                  font_size=12, bold=True, color_hex=theme["text"])
        for i, t in enumerate(trends[:4]):
            ty = 3.55 + i * 0.42
            bar_w = 3.5 - (i * 0.4)  # descending bars for visual
            _add_rect(slide, 0.5, ty, bar_w, 0.28, fill_hex=theme["accent"])
            _add_text(slide, t[:80], 0.5 + bar_w + 0.15, ty, 5.5 - bar_w, 0.3,
                      font_size=9, color_hex=theme["text"])

    # Geographic focus
    geo = _parse_str(mr.get("geography") or mr.get("geographic_focus"))
    if geo:
        _add_rect(slide, 6.0, 2.2, 3.8, 2.85, fill_hex=theme["card"], border_hex=theme["border"])
        _add_text(slide, "Geographic Focus", 6.15, 2.27, 3.5, 0.35,
                  font_size=10, bold=True, color_hex=theme["accent"])
        _add_text(slide, geo[:300], 6.15, 2.65, 3.5, 2.2,
                  font_size=9, color_hex=theme["text"], wrap=True)


def build_audience_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 5, theme)
    _slide_title(slide, "Who We Serve", theme)

    ta = _get_stage(data, "target_audience")
    if not ta:
        return

    personas = ta.get("personas") or ta.get("customer_segments") or []
    if isinstance(personas, dict):
        personas = [personas]

    channels = _parse_list(ta.get("acquisition_channels") or ta.get("channels"), 4)
    audience_size = _parse_str(ta.get("audience_size") or ta.get("market_size"))

    # Persona cards
    for i, p in enumerate(personas[:2]):
        if not isinstance(p, dict):
            continue
        cx = 0.5 + i * 4.8
        _add_rect(slide, cx, 0.75, 4.5, 3.8, fill_hex=theme["card"], border_hex=theme["accent"])
        _add_text(slide, f"Persona {i+1}", cx + 0.2, 0.82, 4.1, 0.35,
                  font_size=9, bold=True, color_hex=theme["accent"])

        fields = [
            ("Role", p.get("role") or p.get("title") or p.get("name")),
            ("Company", p.get("company_size") or p.get("company")),
            ("Pain", p.get("pain") or p.get("pain_point") or p.get("challenge")),
            ("WTP", p.get("willingness_to_pay") or p.get("budget") or p.get("wtp")),
        ]
        fy = 1.2
        for label, val in fields:
            val = _parse_str(val)
            if val:
                _add_text(slide, f"{label}:", cx + 0.2, fy, 1.3, 0.35,
                          font_size=9, bold=True, color_hex=theme["muted"])
                _add_text(slide, val[:60], cx + 1.5, fy, 2.9, 0.35,
                          font_size=9, color_hex=theme["text"])
                fy += 0.42

    # Audience size
    if audience_size:
        _add_rect(slide, 0.5, 4.7, 4.5, 0.55, fill_hex=theme["accent"])
        _add_text(slide, f"Audience Size: {audience_size}", 0.65, 4.77, 4.2, 0.4,
                  font_size=11, bold=True, color_hex="FFFFFF")

    # Acquisition channels
    if channels:
        _add_text(slide, "Acquisition Channels", 5.5, 4.65, 4.0, 0.35,
                  font_size=10, bold=True, color_hex=theme["text"])
        for i, ch in enumerate(channels[:4]):
            _add_rect(slide, 5.5 + i * 1.1, 4.95, 1.0, 0.45, fill_hex=theme["card"], border_hex=theme["accent"])
            _add_text(slide, ch[:12], 5.55 + i * 1.1, 5.0, 0.9, 0.35,
                      font_size=8, color_hex=theme["accent"], align=PP_ALIGN.CENTER)


def build_competitor_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 6, theme)
    _slide_title(slide, "Competitive Landscape", theme)

    ca = _get_stage(data, "competitor_analysis")
    if not ca:
        return

    competitors = ca.get("competitors") or ca.get("competitor_list") or []
    gaps = _parse_list(ca.get("market_gaps") or ca.get("gaps"), 3)
    moat = _parse_str(ca.get("competitive_moat") or ca.get("our_advantage") or ca.get("moat"))

    # Table header
    headers = ["Competitor", "Type", "Threat", "Key Weakness"]
    col_widths = [2.5, 1.5, 1.2, 3.6]
    col_x = [0.5, 3.0, 4.5, 5.7]

    _add_rect(slide, 0.5, 0.75, 9.0, 0.38, fill_hex=theme["accent"])
    for hx, hw, ht in zip(col_x, col_widths, headers):
        _add_text(slide, ht, hx + 0.05, 0.78, hw - 0.1, 0.32,
                  font_size=9, bold=True, color_hex="FFFFFF")

    if isinstance(competitors, list):
        for row_i, comp in enumerate(competitors[:5]):
            if not isinstance(comp, dict):
                continue
            ry = 1.18 + row_i * 0.45
            row_bg = theme["card"] if row_i % 2 == 0 else theme["bg"]
            _add_rect(slide, 0.5, ry, 9.0, 0.42, fill_hex=row_bg, border_hex=theme["border"])

            vals = [
                _parse_str(comp.get("name") or comp.get("competitor"))[:28],
                _parse_str(comp.get("type") or comp.get("category"))[:15],
                _parse_str(comp.get("threat_level") or comp.get("threat"))[:10],
                _parse_str(comp.get("weakness") or comp.get("key_weakness"))[:55],
            ]
            for cx, cw, val in zip(col_x, col_widths, vals):
                if val:
                    _add_text(slide, val, cx + 0.05, ry + 0.05, cw - 0.1, 0.35,
                              font_size=9, color_hex=theme["text"])

    # Market gaps
    if gaps:
        _add_text(slide, "Market Gaps We Exploit", 0.5, 3.85, 5.5, 0.35,
                  font_size=11, bold=True, color_hex=theme["text"])
        for i, g in enumerate(gaps[:3]):
            gx = 0.5 + i * 3.1
            _add_rect(slide, gx, 4.25, 2.9, 0.85, fill_hex=theme["card"], border_hex=theme["accent"])
            _add_text(slide, g[:100], gx + 0.12, 4.32, 2.7, 0.72,
                      font_size=9, color_hex=theme["accent"], wrap=True)

    # Moat
    if moat:
        _add_rect(slide, 0.5, 5.1, 9.0, 0.38, fill_hex=theme["accent"])
        _add_text(slide, f"Our Moat: {moat[:120]}", 0.65, 5.15, 8.7, 0.32,
                  font_size=10, bold=True, color_hex="FFFFFF")


def build_monetization_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 7, theme)
    _slide_title(slide, "Business Model", theme)

    mon = _get_stage(data, "monetization")
    if not mon:
        return

    tiers = mon.get("pricing_tiers") or mon.get("pricing") or []
    metrics = [
        ("MRR Yr 1", mon.get("mrr_year_1") or mon.get("mrr")),
        ("LTV:CAC", mon.get("ltv_cac") or mon.get("ltv_to_cac")),
        ("Payback", mon.get("payback_period") or mon.get("payback")),
    ]

    revenue_model = _parse_str(mon.get("revenue_model") or mon.get("model_description"))
    if revenue_model:
        _add_rect(slide, 0.5, 0.75, 9.0, 0.65, fill_hex=theme["card"], border_hex=theme["border"])
        _add_text(slide, revenue_model[:200], 0.65, 0.82, 8.7, 0.55,
                  font_size=10, color_hex=theme["text"], wrap=True)

    # Pricing tier cards
    if isinstance(tiers, list) and tiers:
        tier_count = min(len(tiers), 4)
        card_w = 9.0 / tier_count - 0.2
        for i, tier in enumerate(tiers[:4]):
            if not isinstance(tier, dict):
                continue
            tx = 0.5 + i * (card_w + 0.2)
            is_featured = tier.get("featured") or i == 2
            border = theme["accent"] if is_featured else theme["border"]
            _add_rect(slide, tx, 1.55, card_w, 2.8, fill_hex=theme["card"], border_hex=border)
            # Plan name
            plan = _parse_str(tier.get("name") or tier.get("plan"))
            price = _parse_str(tier.get("price") or tier.get("price_per_month"))
            features = _parse_list(tier.get("features") or tier.get("includes"), 3)

            if plan:
                _add_text(slide, plan, tx + 0.1, 1.62, card_w - 0.2, 0.38,
                          font_size=11, bold=True, color_hex=theme["accent"] if is_featured else theme["text"])
            if price:
                _add_text(slide, price, tx + 0.1, 1.98, card_w - 0.2, 0.45,
                          font_size=18, bold=True, color_hex=theme["text"])
            for j, feat in enumerate(features[:3]):
                fy = 2.52 + j * 0.42
                _add_text(slide, f"• {feat[:40]}", tx + 0.1, fy, card_w - 0.2, 0.38,
                          font_size=9, color_hex=theme["muted"])

    # Key metrics bar
    has_metrics = any(_parse_str(v) for _, v in metrics)
    if has_metrics:
        _add_rect(slide, 0.5, 4.55, 9.0, 0.7, fill_hex=theme["accent"])
        for i, (label, val) in enumerate(metrics):
            val = _parse_str(val)
            if val:
                mx = 0.9 + i * 3.1
                _add_text(slide, label, mx, 4.6, 2.5, 0.28,
                          font_size=8, bold=True, color_hex="FFFFFF")
                _add_text(slide, val[:25], mx, 4.88, 2.5, 0.32,
                          font_size=13, bold=True, color_hex="FFFFFF")


def build_mvp_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 8, theme)
    _slide_title(slide, "MVP Roadmap", theme)

    mvp = _get_stage(data, "mvp_planning")
    if not mvp:
        return

    phases = mvp.get("phases") or mvp.get("roadmap_phases") or []
    budget = _parse_str(mvp.get("budget") or mvp.get("estimated_budget"))
    team = _parse_str(mvp.get("team_size") or mvp.get("team"))
    timeline = _parse_str(mvp.get("timeline") or mvp.get("total_duration"))

    # 3-phase timeline
    if isinstance(phases, list) and phases:
        for i, phase in enumerate(phases[:3]):
            if not isinstance(phase, dict):
                continue
            px = 0.5 + i * 3.1
            # connector line
            if i < 2:
                _add_rect(slide, px + 2.9, 1.05, 0.2, 0.1, fill_hex=theme["accent"])
            _add_rect(slide, px, 0.75, 2.9, 0.5, fill_hex=theme["accent"])
            phase_name = _parse_str(phase.get("name") or phase.get("phase") or f"Phase {i+1}")
            _add_text(slide, phase_name, px + 0.1, 0.79, 2.7, 0.42,
                      font_size=11, bold=True, color_hex="FFFFFF")

            weeks = _parse_str(phase.get("duration") or phase.get("weeks") or phase.get("timeline"))
            if weeks:
                _add_text(slide, weeks, px + 0.1, 1.3, 2.7, 0.3,
                          font_size=9, color_hex=theme["muted"])

            features = _parse_list(phase.get("features") or phase.get("tasks") or phase.get("deliverables"), 4)
            _add_rect(slide, px, 1.62, 2.9, 2.85, fill_hex=theme["card"], border_hex=theme["border"])
            for j, feat in enumerate(features[:4]):
                fy = 1.72 + j * 0.58
                _add_text(slide, f"• {feat[:55]}", px + 0.12, fy, 2.7, 0.52,
                          font_size=9, color_hex=theme["text"], wrap=True)

    # Bottom metrics
    bottom_items = [(l, v) for l, v in [("Timeline", timeline), ("Budget", budget), ("Team", team)] if v]
    if bottom_items:
        _add_rect(slide, 0.5, 4.6, 9.0, 0.65, fill_hex=theme["card"], border_hex=theme["border"])
        for i, (lbl, val) in enumerate(bottom_items[:3]):
            bx = 0.9 + i * 3.1
            _add_text(slide, lbl, bx, 4.65, 2.5, 0.25, font_size=8, bold=True, color_hex=theme["muted"])
            _add_text(slide, val[:30], bx, 4.9, 2.5, 0.3, font_size=11, bold=True, color_hex=theme["accent"])


def build_gtm_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 9, theme)
    _slide_title(slide, "Go-To-Market Strategy", theme)

    gtm = _get_stage(data, "gtm_strategy")
    if not gtm:
        return

    channels = _parse_list(gtm.get("channels") or gtm.get("acquisition_channels"), 4)
    day1_goal = _parse_str(gtm.get("day_1_goal") or gtm.get("launch_goal"))
    cac = _parse_str(gtm.get("cac_target") or gtm.get("cac"))
    funnel = gtm.get("funnel") or ["Awareness", "Activation", "Conversion", "Revenue"]
    launch_strategy = _parse_str(gtm.get("launch_strategy") or gtm.get("strategy"))

    # Channel cards
    for i, ch in enumerate(channels[:4]):
        cx = 0.5 + i * 2.35
        _add_rect(slide, cx, 0.75, 2.15, 1.6, fill_hex=theme["card"], border_hex=theme["accent"])
        _add_text(slide, ch[:60], cx + 0.12, 0.85, 1.95, 1.4,
                  font_size=9, color_hex=theme["text"], wrap=True)

    # Funnel
    if funnel:
        funnel_items = funnel if isinstance(funnel, list) else [funnel]
        _add_text(slide, "Launch Funnel", 0.5, 2.55, 4.0, 0.35,
                  font_size=11, bold=True, color_hex=theme["text"])
        funnel_colors = [theme["accent"], theme["accent2"], theme["success"], theme["warning"]]
        for i, step in enumerate(funnel_items[:4]):
            step = _parse_str(step)
            sw = 4.2 - i * 0.5
            sx = 0.5 + (i * 0.25)
            _add_rect(slide, sx, 3.0 + i * 0.55, sw, 0.45, fill_hex=funnel_colors[i % 4])
            _add_text(slide, f"→ {step[:35]}", sx + 0.15, 3.06 + i * 0.55, sw - 0.3, 0.36,
                      font_size=10, bold=True, color_hex="FFFFFF")

    # Key metrics
    if day1_goal:
        _add_rect(slide, 5.5, 2.55, 4.2, 0.7, fill_hex=theme["accent"])
        _add_text(slide, "Day-1 Goal", 5.65, 2.58, 3.9, 0.28,
                  font_size=8, bold=True, color_hex="FFFFFF")
        _add_text(slide, day1_goal[:80], 5.65, 2.86, 3.9, 0.36,
                  font_size=10, bold=True, color_hex="FFFFFF", wrap=True)

    if cac:
        _add_rect(slide, 5.5, 3.35, 4.2, 0.7, fill_hex=theme["card"], border_hex=theme["accent"])
        _add_text(slide, "CAC Target", 5.65, 3.38, 3.9, 0.28,
                  font_size=8, bold=True, color_hex=theme["muted"])
        _add_text(slide, cac[:40], 5.65, 3.66, 3.9, 0.36,
                  font_size=14, bold=True, color_hex=theme["accent"])

    if launch_strategy:
        _add_rect(slide, 5.5, 4.15, 4.2, 1.1, fill_hex=theme["card"], border_hex=theme["border"])
        _add_text(slide, launch_strategy[:200], 5.65, 4.22, 4.0, 1.0,
                  font_size=9, color_hex=theme["text"], wrap=True)


def build_closing_slide(slide, theme: dict, data: dict):
    _set_slide_bg(slide, theme["cover_bg"])
    _add_top_bar(slide, theme["accent"])
    _add_slide_number(slide, 10, theme)
    _slide_title(slide, "Why Now. Why Us.", theme)

    fr = _get_stage(data, "final_report")
    score = data.get("overall_score") or 0
    risk = _parse_str(fr.get("risk_level") if fr else None)
    recommendations = _parse_list(fr.get("recommendations") or fr.get("top_recommendations") if fr else None, 3)

    # Big score display
    if score and score > 0:
        score_col = _score_color(score, data.get("_theme", theme))
        _add_rect(slide, 0.5, 0.85, 1.6, 1.5, fill_hex=score_col)
        _add_text(slide, str(score), 0.55, 0.95, 1.5, 1.2,
                  font_size=52, bold=True, color_hex="FFFFFF", align=PP_ALIGN.CENTER)
        _add_text(slide, "/ 100", 0.55, 2.05, 1.5, 0.35,
                  font_size=10, color_hex="FFFFFF", align=PP_ALIGN.CENTER)

    # Risk badge
    if risk:
        risk_colors = {"low": theme["success"], "medium": theme["warning"], "high": theme["danger"]}
        rc = risk_colors.get(risk.lower(), theme["muted"])
        _add_rect(slide, 2.3, 1.05, 1.8, 0.55, fill_hex=rc)
        _add_text(slide, f"Risk: {risk.title()}", 2.45, 1.1, 1.55, 0.45,
                  font_size=11, bold=True, color_hex="FFFFFF")

    # Top recommendations
    if recommendations:
        _add_text(slide, "Strategic Recommendations", 0.5, 2.6, 6.0, 0.38,
                  font_size=13, bold=True, color_hex=theme["text"])
        for i, rec in enumerate(recommendations[:3]):
            ry = 3.1 + i * 0.72
            _add_rect(slide, 0.5, ry, 9.0, 0.62, fill_hex=theme["card"], border_hex=theme["accent"])
            _add_rect(slide, 0.5, ry, 0.06, 0.62, fill_hex=theme["accent"])
            _add_text(slide, f"{i+1}.", 0.65, ry + 0.1, 0.4, 0.42,
                      font_size=12, bold=True, color_hex=theme["accent"])
            _add_text(slide, rec[:180], 1.12, ry + 0.1, 8.25, 0.48,
                      font_size=10, color_hex=theme["text"], wrap=True)

    # Brand tagline
    _add_text(slide,
              "Inceptrax — Build less. Validate more. Execute with confidence.",
              0.5, 5.1, 9.0, 0.3,
              font_size=9, color_hex=theme["muted"], align=PP_ALIGN.CENTER)


# ─── Main Generator ────────────────────────────────────────────────────────────

SLIDE_BUILDERS = [
    build_cover_slide,
    build_problem_slide,
    build_solution_slide,
    build_market_slide,
    build_audience_slide,
    build_competitor_slide,
    build_monetization_slide,
    build_mvp_slide,
    build_gtm_slide,
    build_closing_slide,
]


def generate_investor_ppt(analysis_data: dict, theme_key: str = "dark_executive") -> str:
    """
    Generate a 10-slide investor deck with the chosen theme.
    Only renders real user data — missing fields are silently skipped.

    Args:
        analysis_data: dict with idea info + stages dict
        theme_key: one of 'dark_executive', 'clean_light', 'gradient_bold'

    Returns:
        Absolute path to the generated .pptx file
    """
    theme = THEMES.get(theme_key, THEMES["dark_executive"])
    analysis_data["_theme"] = theme  # pass for closing slide score color

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    blank_layout = prs.slide_layouts[6]  # completely blank

    for builder in SLIDE_BUILDERS:
        slide = prs.slides.add_slide(blank_layout)
        try:
            builder(slide, theme, analysis_data)
        except Exception as e:
            # Never crash the whole deck over one bad slide
            print(f"[PPT] Slide build error in {builder.__name__}: {e}")

    idea_id = analysis_data.get("id", "unknown")
    out_dir = tempfile.gettempdir()
    out_path = os.path.join(out_dir, f"inceptrax_{idea_id}_{theme_key}.pptx")
    prs.save(out_path)
    return out_path


# ─── Legacy compatibility shim ────────────────────────────────────────────────

class PPTService:
    """Backwards-compatible shim for existing route that calls PPTService.generate_presentation()."""

    @staticmethod
    def generate_presentation(idea) -> str:
        """Build analysis_data from idea ORM object and call generate_investor_ppt."""
        from app.models.user_model import StageResult
        import json

        stages = {}
        for sr in (idea.stage_results or []):
            try:
                stages[sr.stage_name] = json.loads(sr.result_json or "{}")
            except Exception:
                stages[sr.stage_name] = {}

        data = {
            "id": idea.id,
            "title": idea.title,
            "description": idea.description,
            "one_liner": idea.one_liner if hasattr(idea, "one_liner") else "",
            "industry": idea.industry or idea.market or "",
            "overall_score": idea.overall_score or 0,
            "stages": stages,
        }
        return generate_investor_ppt(data, "dark_executive")
